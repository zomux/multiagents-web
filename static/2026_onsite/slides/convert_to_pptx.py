#!/usr/bin/env python3
"""
Convert WMAC 2026 slides to PowerPoint format.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors from the HTML slides
NAVY = RGBColor(0x1a, 0x36, 0x5d)  # #1a365d
TEAL = RGBColor(0x31, 0x97, 0x95)  # #319795
GRAY = RGBColor(0x71, 0x80, 0x96)  # #718096
LIGHT_BG = RGBColor(0xf7, 0xfa, 0xfc)  # #f7fafc
WHITE = RGBColor(0xff, 0xff, 0xff)

# Slide dimensions (16:9 standard)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=NAVY):
    """Add a bullet list to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.level = 0
    return txBox

def add_image_safe(slide, img_path, left, top, width=None, height=None):
    """Add an image to the slide if it exists."""
    if os.path.exists(img_path):
        try:
            if width and height:
                return slide.shapes.add_picture(img_path, left, top, width, height)
            elif width:
                return slide.shapes.add_picture(img_path, left, top, width=width)
            elif height:
                return slide.shapes.add_picture(img_path, left, top, height=height)
            else:
                return slide.shapes.add_picture(img_path, left, top)
        except Exception as e:
            print(f"Warning: Could not add image {img_path}: {e}")
    else:
        print(f"Warning: Image not found: {img_path}")
    return None

def create_title_slide(prs, slide_layout):
    """Slide 1: Title slide"""
    slide = prs.slides.add_slide(slide_layout)

    # Add logo
    logo_path = "logo_wmac.png"
    add_image_safe(slide, logo_path, Inches(2.5), Inches(1.5), width=Inches(1.5))

    # Title
    add_text_box(slide, Inches(4.2), Inches(1.5), Inches(5), Inches(0.8),
                 "WMAC 2026", font_size=48, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    # Subtitle
    add_text_box(slide, Inches(4.2), Inches(2.2), Inches(5), Inches(0.4),
                 "AAAI 2026 Bridge Program on", font_size=14, bold=True, color=NAVY)
    add_text_box(slide, Inches(4.2), Inches(2.5), Inches(5), Inches(0.4),
                 "Advancing LLM-Based Multi-Agent Collaboration", font_size=18, bold=True, color=NAVY)

    # Date
    add_text_box(slide, Inches(4.2), Inches(3.0), Inches(5), Inches(0.3),
                 "January 20, 2026, Singapore", font_size=12, color=GRAY)

    # Website
    add_text_box(slide, Inches(0.5), Inches(5.0), Inches(3), Inches(0.3),
                 "multiagents.org", font_size=12, color=TEAL)

def create_organizers_slide(prs, slide_layout):
    """Slide 2: Overview & Organizers"""
    slide = prs.slides.add_slide(slide_layout)

    # Title
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
                 "Advancing LLM-Based Multi-Agent Collaboration", font_size=18, bold=True, color=NAVY)

    # Overview bullets
    overview_items = [
        "AI Agents driven by LLMs are becoming capable of complex task solving",
        "Rise of Multi-Agent Collaboration",
        "Open challenges remain"
    ]
    add_bullet_list(slide, Inches(0.5), Inches(0.75), Inches(9), Inches(0.8), overview_items, font_size=11, color=NAVY)

    # Section title
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.35),
                 "Program Organizers", font_size=13, bold=True, color=NAVY)

    # Organizer info (text only, images would need to be added manually)
    organizers_row1 = [
        ("Raphael Shu", "Chair"),
        ("Yusen Zhang", "Co-Chair"),
        ("Alborz Geramifard", "LinkedIn"),
        ("Alex Marin", "Thomson Reuters")
    ]
    organizers_row2 = [
        ("Weiyan Shi", "Northeastern"),
        ("Tao Yu", "HKU"),
        ("Yi Zhang", "AWS"),
        ("Jeffrey Cho", "UPenn")
    ]

    # Add organizer photos (Row 1)
    photo_files_row1 = [
        "photos/circle_raphael_shu.jpeg",
        "photos/circle_yusen_zhang.png",
        "photos/circle_alborz_geramifard.jpg",
        "photos/circle_alex_marin.png"
    ]
    photo_files_row2 = [
        "photos/circle_weiyan_shi.jpg",
        "photos/circle_tao_yu.jpeg",
        "photos/circle_yi_zhang.png",
        "photos/circle_jeffrey_cho.jpg"
    ]

    for i, ((name, role), photo) in enumerate(zip(organizers_row1, photo_files_row1)):
        x_pos = Inches(1.2 + i * 2.2)
        add_image_safe(slide, photo, x_pos, Inches(1.85), width=Inches(0.9), height=Inches(0.9))
        add_text_box(slide, x_pos - Inches(0.3), Inches(2.8), Inches(1.5), Inches(0.25),
                     name, font_size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text_box(slide, x_pos - Inches(0.3), Inches(3.0), Inches(1.5), Inches(0.2),
                     role, font_size=7, color=TEAL if role in ["Chair", "Co-Chair"] else GRAY, align=PP_ALIGN.CENTER)

    for i, ((name, role), photo) in enumerate(zip(organizers_row2, photo_files_row2)):
        x_pos = Inches(1.2 + i * 2.2)
        add_image_safe(slide, photo, x_pos, Inches(3.35), width=Inches(0.9), height=Inches(0.9))
        add_text_box(slide, x_pos - Inches(0.3), Inches(4.3), Inches(1.5), Inches(0.25),
                     name, font_size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text_box(slide, x_pos - Inches(0.3), Inches(4.5), Inches(1.5), Inches(0.2),
                     role, font_size=7, color=GRAY, align=PP_ALIGN.CENTER)

def create_sponsors_slide(prs, slide_layout):
    """Slide 3: Sponsors"""
    slide = prs.slides.add_slide(slide_layout)

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.6),
                 "SPONSORS", font_size=24, bold=True, color=NAVY)
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(9), Inches(0.4),
                 "We thank our sponsors for their generous support", font_size=12, color=GRAY)

    # Sponsor logos
    sponsor_logos = [
        "../../2026_sponsors/linkedin_logo.png",
        "../../2026_sponsors/aws_logo.png",
        "../../2026_sponsors/youware_logo.png"
    ]

    for i, logo in enumerate(sponsor_logos):
        add_image_safe(slide, logo, Inches(1.5 + i * 2.8), Inches(2.2), height=Inches(0.8))

    add_text_box(slide, Inches(0.5), Inches(5.0), Inches(3), Inches(0.3),
                 "multiagents.org/2026", font_size=12, color=TEAL)

def create_agenda_slide(prs, slide_layout):
    """Slide 4: Agenda"""
    slide = prs.slides.add_slide(slide_layout)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
                 "WMAC 2026 PROGRAM AGENDA", font_size=24, bold=True, color=NAVY)

    # Morning sessions
    add_text_box(slide, Inches(0.5), Inches(0.9), Inches(4), Inches(0.4),
                 "MORNING SESSIONS 9:00 ~ 12:30", font_size=12, bold=True, color=TEAL)
    morning_items = [
        "Opening",
        "Invited Talk 1 (Yi R. Fung)",
        "WMAC History and Vision",
        "Lightning Talks for Posters (30 mins)",
        "Coffee Break",
        "Invited Talk 2 (Sophia Han)",
        "Poster Session"
    ]
    for i, item in enumerate(morning_items):
        add_text_box(slide, Inches(0.5), Inches(1.3 + i * 0.35), Inches(4), Inches(0.35),
                     item, font_size=11, color=NAVY)

    # Afternoon sessions
    add_text_box(slide, Inches(5.2), Inches(0.9), Inches(4.5), Inches(0.4),
                 "AFTERNOON SESSIONS 13:50 ~ 18:15", font_size=12, bold=True, color=TEAL)
    afternoon_items = [
        "Invited Talk 3 (Matthew Taylor)",
        "Specially Invited Oral Talk",
        "Fireside Chat & Panel",
        "Oral Presentations 2 & 3",
        "Industrial Demo Session",
        "Oral Presentation 4",
        "Award & Closing"
    ]
    for i, item in enumerate(afternoon_items):
        add_text_box(slide, Inches(5.2), Inches(1.3 + i * 0.35), Inches(4.5), Inches(0.35),
                     item, font_size=11, color=NAVY)

    add_text_box(slide, Inches(0.5), Inches(5.0), Inches(3), Inches(0.3),
                 "multiagents.org/2026", font_size=12, color=TEAL)

def create_attendance_slide(prs, slide_layout):
    """Slide 5: Attendance Resources"""
    slide = prs.slides.add_slide(slide_layout)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
                 "ATTENDANCE RESOURCES", font_size=24, bold=True, color=NAVY)

    sections = [
        ("Room: Topaz 220 – 225", "See venue map on multiagents.org/2026"),
        ("Registration", "Register with \"AAAI Tutorial/Lab/Bridge only\" or\n\"Bridges, Tutorials and Labs Add-on\" option"),
        ("Papers & Posters", "Available at multiagents.org/2026"),
        ("Remote Participation", "In-person encouraged, remote supported for talks\nmultiagents.org/2026_remote")
    ]

    y_pos = 0.9
    for title, desc in sections:
        add_text_box(slide, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.35),
                     title, font_size=14, bold=True, color=NAVY)
        add_text_box(slide, Inches(0.5), Inches(y_pos + 0.35), Inches(9), Inches(0.6),
                     desc, font_size=12, color=GRAY)
        y_pos += 1.0

    add_text_box(slide, Inches(0.5), Inches(5.0), Inches(3), Inches(0.3),
                 "multiagents.org/2026", font_size=12, color=TEAL)

def create_questions_slide(prs, slide_layout):
    """Slide 6: Questions"""
    slide = prs.slides.add_slide(slide_layout)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
                 "HOW TO ASK QUESTIONS", font_size=24, bold=True, color=NAVY)
    add_text_box(slide, Inches(0.5), Inches(0.85), Inches(9), Inches(0.4),
                 "You can raise questions to speakers or the event organizers through:", font_size=12, color=GRAY)

    # Discord
    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4),
                 "Discord", font_size=16, bold=True, color=NAVY)
    add_text_box(slide, Inches(0.5), Inches(1.9), Inches(9), Inches(0.3),
                 "Join our Discord server for real-time discussion", font_size=12, color=GRAY)
    add_text_box(slide, Inches(0.5), Inches(2.2), Inches(9), Inches(0.3),
                 "https://discord.gg/PLACEHOLDER", font_size=12, color=TEAL)

    # Underline
    add_text_box(slide, Inches(0.5), Inches(2.8), Inches(9), Inches(0.4),
                 "Underline", font_size=16, bold=True, color=NAVY)
    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(9), Inches(0.3),
                 "Submit questions through the Underline platform", font_size=12, color=GRAY)
    add_text_box(slide, Inches(0.5), Inches(3.5), Inches(9), Inches(0.3),
                 "https://underline.io/PLACEHOLDER", font_size=12, color=TEAL)

    add_text_box(slide, Inches(0.5), Inches(5.0), Inches(3), Inches(0.3),
                 "multiagents.org/2026", font_size=12, color=TEAL)

def create_reminders_slide(prs, slide_layout):
    """Slide 7: Reminders"""
    slide = prs.slides.add_slide(slide_layout)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
                 "REMINDERS", font_size=24, bold=True, color=NAVY)

    # Goodies
    add_text_box(slide, Inches(0.5), Inches(1.0), Inches(9), Inches(0.4),
                 "Goodies", font_size=16, bold=True, color=NAVY)
    add_text_box(slide, Inches(0.5), Inches(1.4), Inches(9), Inches(0.6),
                 "We have prepared t-shirts and tote bags for attendees.\nPlease check in with a volunteer to obtain your goodies.", font_size=12, color=GRAY)

    # Lunchtime Gathering
    add_text_box(slide, Inches(0.5), Inches(2.3), Inches(9), Inches(0.4),
                 "Lunchtime Gathering", font_size=16, bold=True, color=NAVY)
    add_text_box(slide, Inches(0.5), Inches(2.7), Inches(9), Inches(0.6),
                 "After the poster session, please come back to this room.\nWe may have lunchbox and drinks available for attendees.", font_size=12, color=GRAY)

    add_text_box(slide, Inches(0.5), Inches(5.0), Inches(3), Inches(0.3),
                 "multiagents.org/2026", font_size=12, color=TEAL)

def create_presenters_slide(prs, slide_layout):
    """Slide 8: For Presenters"""
    slide = prs.slides.add_slide(slide_layout)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
                 "FOR PRESENTERS", font_size=24, bold=True, color=NAVY)

    items = [
        "Confirm program agenda",
        "Check in with moderators 15 mins before presentation",
        "All presentations will be in-person",
        "Oral presentations: 20 mins each"
    ]
    add_bullet_list(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(2.5), items, font_size=14, color=NAVY)

    add_text_box(slide, Inches(0.5), Inches(5.0), Inches(3), Inches(0.3),
                 "multiagents.org/2026", font_size=12, color=TEAL)

def create_blank_slide(prs, slide_layout):
    """Slide 9: Blank"""
    slide = prs.slides.add_slide(slide_layout)
    # Just a blank slide

def create_event_slide(prs, slide_layout, event_title, event_time=None):
    """Create a centered event slide (Coffee Break, Poster Session, etc.)"""
    slide = prs.slides.add_slide(slide_layout)

    # Add logo
    logo_path = "logo_wmac.png"
    add_image_safe(slide, logo_path, Inches(2.5), Inches(0.8), width=Inches(1.2))

    # Title
    add_text_box(slide, Inches(3.9), Inches(0.8), Inches(5.5), Inches(0.7),
                 "WMAC 2026", font_size=40, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    add_text_box(slide, Inches(3.9), Inches(1.4), Inches(5.5), Inches(0.3),
                 "AAAI 2026 Bridge Program on", font_size=12, bold=True, color=NAVY)
    add_text_box(slide, Inches(3.9), Inches(1.65), Inches(5.5), Inches(0.35),
                 "Advancing LLM-Based Multi-Agent Collaboration", font_size=14, bold=True, color=NAVY)
    add_text_box(slide, Inches(3.9), Inches(1.95), Inches(5.5), Inches(0.25),
                 "January 20, 2026, Singapore", font_size=10, color=GRAY)

    # Event title
    add_text_box(slide, Inches(0.5), Inches(3.0), Inches(9), Inches(0.6),
                 event_title, font_size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    if event_time:
        add_text_box(slide, Inches(0.5), Inches(3.6), Inches(9), Inches(0.4),
                     event_time, font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(5.0), Inches(3), Inches(0.3),
                 "multiagents.org", font_size=12, color=TEAL)

def create_best_paper_slide(prs, slide_layout):
    """Slide: Best Paper Award"""
    slide = prs.slides.add_slide(slide_layout)

    # Add logo
    logo_path = "logo_wmac.png"
    add_image_safe(slide, logo_path, Inches(0.5), Inches(0.3), width=Inches(1.0))

    add_text_box(slide, Inches(1.7), Inches(0.3), Inches(7), Inches(0.6),
                 "WMAC 2026", font_size=36, bold=True, color=NAVY)
    add_text_box(slide, Inches(1.7), Inches(0.85), Inches(7.5), Inches(0.3),
                 "AAAI 2026 Bridge Program on Advancing LLM-Based Multi-Agent Collaboration", font_size=11, bold=True, color=NAVY)

    add_text_box(slide, Inches(0.5), Inches(2.5), Inches(9), Inches(0.8),
                 "BEST PAPER AWARD", font_size=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

def create_closing_slide(prs, slide_layout):
    """Slide: Closing Remarks"""
    slide = prs.slides.add_slide(slide_layout)

    # Add logo
    logo_path = "logo_wmac.png"
    add_image_safe(slide, logo_path, Inches(0.5), Inches(0.3), width=Inches(1.0))

    add_text_box(slide, Inches(1.7), Inches(0.3), Inches(7), Inches(0.6),
                 "WMAC 2026", font_size=36, bold=True, color=NAVY)
    add_text_box(slide, Inches(1.7), Inches(0.85), Inches(7.5), Inches(0.3),
                 "AAAI 2026 Bridge Program on Advancing LLM-Based Multi-Agent Collaboration", font_size=11, bold=True, color=NAVY)

    add_text_box(slide, Inches(0.5), Inches(2.5), Inches(9), Inches(0.8),
                 "CLOSING REMARKS", font_size=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

def create_thank_you_slide(prs, slide_layout):
    """Slide: Thank You"""
    slide = prs.slides.add_slide(slide_layout)

    # Add logo
    logo_path = "logo_wmac.png"
    add_image_safe(slide, logo_path, Inches(0.5), Inches(0.3), width=Inches(1.0))

    add_text_box(slide, Inches(1.7), Inches(0.3), Inches(7), Inches(0.6),
                 "WMAC 2026", font_size=36, bold=True, color=NAVY)
    add_text_box(slide, Inches(1.7), Inches(0.85), Inches(7.5), Inches(0.3),
                 "AAAI 2026 Bridge Program on Advancing LLM-Based Multi-Agent Collaboration", font_size=11, bold=True, color=NAVY)

    add_text_box(slide, Inches(0.5), Inches(2.5), Inches(9), Inches(0.8),
                 "THANK YOU", font_size=40, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # Create all slides
    print("Creating slides...")

    # Slide 1: Title
    create_title_slide(prs, blank_layout)
    print("  1. Title")

    # Slide 2: Organizers
    create_organizers_slide(prs, blank_layout)
    print("  2. Overview & Organizers")

    # Slide 3: Sponsors
    create_sponsors_slide(prs, blank_layout)
    print("  3. Sponsors")

    # Slide 4: Agenda
    create_agenda_slide(prs, blank_layout)
    print("  4. Agenda")

    # Slide 5: Attendance Resources
    create_attendance_slide(prs, blank_layout)
    print("  5. Attendance Resources")

    # Slide 6: Questions
    create_questions_slide(prs, blank_layout)
    print("  6. Questions")

    # Slide 7: Reminders
    create_reminders_slide(prs, blank_layout)
    print("  7. Reminders")

    # Slide 8: For Presenters
    create_presenters_slide(prs, blank_layout)
    print("  8. For Presenters")

    # Slide 9: Blank
    create_blank_slide(prs, blank_layout)
    print("  9. Blank")

    # Slide 10: Coffee Break 10:45
    create_event_slide(prs, blank_layout, "COFFEE BREAK UNTIL 10:45")
    print("  10. Coffee Break 10:45")

    # Slide 11: Poster Session
    create_event_slide(prs, blank_layout, "POSTER SESSION 11:30 - 12:30")
    print("  11. Poster Session")

    # Slide 12: Lunch Break
    create_event_slide(prs, blank_layout, "LUNCH BREAK 12:30 - 13:50")
    print("  12. Lunch Break")

    # Slide 13: Coffee Break 16:05
    create_event_slide(prs, blank_layout, "COFFEE BREAK UNTIL 16:05")
    print("  13. Coffee Break 16:05")

    # Slide 14: Coffee Break 17:00
    create_event_slide(prs, blank_layout, "COFFEE BREAK UNTIL 17:00")
    print("  14. Coffee Break 17:00")

    # Slide 15: Fireside Chat & Panel
    create_event_slide(prs, blank_layout, "FIRESIDE CHAT & PANEL")
    print("  15. Fireside Chat & Panel")

    # Slide 16: Industrial Demo Session
    create_event_slide(prs, blank_layout, "INDUSTRIAL DEMO SESSION", "17:00 - 17:45 (3 demos, 15 mins each)")
    print("  16. Industrial Demo Session")

    # Slide 17: WMAC Title (repeat)
    create_title_slide(prs, blank_layout)
    print("  17. WMAC Title")

    # Slide 18: Best Paper Award
    create_best_paper_slide(prs, blank_layout)
    print("  18. Best Paper Award")

    # Slide 19: Closing Remarks
    create_closing_slide(prs, blank_layout)
    print("  19. Closing Remarks")

    # Slide 20: Thank You
    create_thank_you_slide(prs, blank_layout)
    print("  20. Thank You")

    # Save presentation
    output_file = "slides_combined.pptx"
    prs.save(output_file)
    print(f"\nSaved: {output_file}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
